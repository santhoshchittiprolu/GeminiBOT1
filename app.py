import streamlit as st
import google.generativeai as genai
import os
import fitz  # PyMuPDF for PDFs
import docx
from pptx import Presentation

# Configure Gemini
GOOGLE_API_KEY = os.getenv("AIzaSyDzboVyiU1vAuXkfzF6C5XsUMj1E2M1eDM")
if not GOOGLE_API_KEY:
    st.error("Please set the GEMINI_API_KEY environment variable.")
    st.stop()

genai.configure(api_key=GOOGLE_API_KEY)

# Initialize the model
model = genai.GenerativeModel('gemini-1.5-pro-latest')

# Initialize session state
if "chat" not in st.session_state:
    st.session_state.chat = model.start_chat(history=[])
if "uploaded_files" not in st.session_state:
    st.session_state.uploaded_files = []
if "file_texts" not in st.session_state:
    st.session_state.file_texts = ""

st.title("📚 Lecture Files Chat Assistant")
st.markdown("Upload your lecture files and chat with them using Gemini 1.5 Pro!")

# Function to extract text from files
def extract_text(file):
    if file.type == "application/pdf":
        doc = fitz.open(stream=file.read(), filetype="pdf")
        return "\n".join([page.get_text("text") for page in doc])
    elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":  # DOCX
        doc = docx.Document(file)
        return "\n".join([para.text for para in doc.paragraphs])
    elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":  # PPTX
        prs = Presentation(file)
        return "\n".join([slide.shapes.title.text for slide in prs.slides if slide.shapes.title])
    elif file.type == "text/plain":
        return file.read().decode("utf-8")
    return ""

# File upload section
with st.sidebar:
    st.header("📁 Upload Files")
    uploaded_files = st.file_uploader(
        "Choose lecture files (PDF, TXT, PPTX, DOCX)",
        type=['pdf', 'txt', 'pptx', 'docx'],
        accept_multiple_files=True
    )
    
    if uploaded_files:
        st.session_state.file_texts = ""
        for file in uploaded_files:
            extracted_text = extract_text(file)
            st.session_state.file_texts += f"\n\n--- {file.name} ---\n" + extracted_text
        st.success(f"✅ Uploaded {len(uploaded_files)} files!")

    if st.session_state.file_texts:
        st.subheader("Uploaded Files Content Preview")
        st.text_area("Extracted Text", st.session_state.file_texts[:1000] + "...", height=200)

# Chat interface
for message in st.session_state.chat.history:
    with st.chat_message(message.role):
        st.markdown(message.parts[0].text if message.parts else "")

user_input = st.chat_input("Ask something about your lecture files...")

if user_input:
    with st.chat_message("user"):
        st.markdown(user_input)

    # Construct query with file contents
    full_query = f"Lecture Notes:\n{st.session_state.file_texts}\n\nUser Question: {user_input}"
    
    # Generate response
    try:
        response = st.session_state.chat.send_message(full_query)
        with st.chat_message("assistant"):
            st.markdown(response.candidates[0].text if response.candidates else "(No response)")
    except Exception as e:
        st.error(f"🚨 Error generating response: {str(e)}")

# Add some styling
st.markdown("""
<style>
    .stChatInput textarea {
        min-height: 100px !important;
    }
    .stDownloadButton button {
        width: 100%;
    }
</style>
""", unsafe_allow_html=True)
